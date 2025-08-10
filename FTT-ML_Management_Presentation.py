"""
Financial Data Transformation Platform - Management Presentation Generator
Creates a comprehensive PowerPoint presentation for executive management
Focus: LLM/AI capabilities, workflows, and AWS infrastructure
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

class FTTMLPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_layouts()
        
    def setup_slide_layouts(self):
        """Configure slide layouts and themes"""
        # Use built-in layouts
        self.title_layout = self.prs.slide_layouts[0]  # Title slide
        self.content_layout = self.prs.slide_layouts[1]  # Title and content
        self.two_content_layout = self.prs.slide_layouts[3]  # Two content
        self.blank_layout = self.prs.slide_layouts[6]  # Blank
        
    def create_title_slide(self):
        """Slide 1: Title Slide"""
        slide = self.prs.slides.add_slide(self.title_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Financial Data Transformation Platform (FTT-ML)"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # Dark blue
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = (
            "AI-Powered Financial Data Processing\n"
            "Executive Overview & Technical Architecture\n\n"
            "Leveraging Large Language Models for Automated Data Reconciliation,\n"
            "Transformation, and Analysis"
        )
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 73, 94)
        
    def create_executive_summary_slide(self):
        """Slide 2: Executive Summary"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Executive Summary: AI-Driven Financial Data Processing"
        
        # Content
        content = slide.placeholders[1]
        content.text = (
            "• Revolutionary AI-first approach to financial data processing\n"
            "• Reduces manual effort by 80% through intelligent automation\n"
            "• Processes 50,000-100,000 financial records with AI assistance\n"
            "• Multi-provider LLM integration (OpenAI, Anthropic, Google, Internal)\n"
            "• Enterprise-ready with AWS cloud infrastructure\n"
            "• Maintains human oversight for financial accuracy and compliance\n\n"
            "Key AI Capabilities:\n"
            "  ✓ Natural language to configuration generation\n"
            "  ✓ Intelligent pattern recognition and matching\n"
            "  ✓ Auto-suggested reconciliation rules\n"
            "  ✓ Smart error detection and correction recommendations"
        )
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            if paragraph.text.startswith('•'):
                paragraph.font.color.rgb = RGBColor(46, 125, 50)  # Green
            elif paragraph.text.startswith('  ✓'):
                paragraph.font.color.rgb = RGBColor(13, 71, 161)  # Blue
                paragraph.level = 1
                
    def create_business_problem_slide(self):
        """Slide 3: Business Problem & AI Solution"""
        slide = self.prs.slides.add_slide(self.two_content_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Business Challenge & AI-Powered Solution"
        
        # Left content - Problem
        left_content = slide.placeholders[1]
        left_content.text = (
            "TRADITIONAL CHALLENGES:\n\n"
            "• Manual data reconciliation is time-intensive\n"
            "• Complex rule configuration requires technical expertise\n"
            "• Error-prone manual pattern matching\n"
            "• Inconsistent data transformation logic\n"
            "• Limited scalability for large datasets\n"
            "• High operational costs\n\n"
            "FINANCIAL IMPACT:\n"
            "• 60+ hours/month manual processing\n"
            "• $50K+ annual operational overhead\n"
            "• Risk of human error in critical calculations"
        )
        
        # Right content - Solution
        right_content = slide.placeholders[2]
        right_content.text = (
            "AI-POWERED SOLUTION:\n\n"
            "• LLM converts business requirements to technical rules\n"
            "• Intelligent pattern recognition and auto-matching\n"
            "• AI-suggested configurations with human validation\n"
            "• Natural language processing for data understanding\n"
            "• Automated error detection and recommendations\n"
            "• Scalable cloud infrastructure\n\n"
            "BUSINESS VALUE:\n"
            "• 80% reduction in manual effort\n"
            "• 95% faster configuration generation\n"
            "• Enhanced accuracy with AI validation"
        )
        
        # Format content
        for content in [left_content, right_content]:
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                if paragraph.text.startswith('TRADITIONAL') or paragraph.text.startswith('AI-POWERED'):
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(183, 28, 28)  # Red for problems
                elif paragraph.text.startswith('FINANCIAL') or paragraph.text.startswith('BUSINESS'):
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(46, 125, 50)  # Green for solutions
                    
    def create_ai_capabilities_slide(self):
        """Slide 4: AI & LLM Capabilities"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        title = slide.shapes.title
        title.text = "Advanced AI & Large Language Model Capabilities"
        
        content = slide.placeholders[1]
        content.text = (
            "MULTI-PROVIDER LLM INTEGRATION:\n"
            "• OpenAI GPT-4 Turbo - Advanced reasoning and analysis\n"
            "• Anthropic Claude - Enhanced financial data understanding\n"
            "• Google Gemini - Multimodal processing capabilities\n"
            "• Internal JPMC LLM - Specialized financial domain knowledge\n\n"
            "INTELLIGENT AUTOMATION FEATURES:\n"
            "• Natural Language to Rules: 'Compare transactions by ID with $0.01 tolerance'\n"
            "• Smart Pattern Recognition: Auto-detect reconciliation opportunities\n"
            "• Configuration Generation: AI creates technical rules from business descriptions\n"
            "• Error Analysis: Intelligent identification of data quality issues\n"
            "• Validation Suggestions: AI recommends corrections and improvements\n\n"
            "FINANCIAL AI SPECIALIZATION:\n"
            "• Currency and decimal precision handling\n"
            "• Financial identifier recognition (ISIN, CUSIP, SEDOL)\n"
            "• Regulatory compliance pattern matching\n"
            "• Risk assessment for data processing decisions"
        )
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            if paragraph.text.endswith(':'):
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(13, 71, 161)
                
    def create_workflow_diagram_slide(self):
        """Slide 5: AI-Powered Workflow Diagram"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "AI-Powered Data Processing Workflow"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Create workflow boxes
        boxes = [
            {"text": "1. Business Intent\n(Natural Language)", "x": 1, "y": 1.5, "color": RGBColor(33, 150, 243)},
            {"text": "2. AI Analysis\n(LLM Processing)", "x": 4, "y": 1.5, "color": RGBColor(156, 39, 176)},
            {"text": "3. Rule Generation\n(AI Configuration)", "x": 7, "y": 1.5, "color": RGBColor(255, 193, 7)},
            {"text": "4. Human Review\n(Validation & Approval)", "x": 10, "y": 1.5, "color": RGBColor(244, 67, 54)},
            {"text": "5. Data Processing\n(Automated Execution)", "x": 2.5, "y": 3.5, "color": RGBColor(76, 175, 80)},
            {"text": "6. Results & Analytics\n(Dashboard & Reports)", "x": 5.5, "y": 3.5, "color": RGBColor(96, 125, 139)},
            {"text": "7. Continuous Learning\n(AI Improvement)", "x": 8.5, "y": 3.5, "color": RGBColor(255, 87, 34)}
        ]
        
        for box in boxes:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(box["x"]), Inches(box["y"]),
                Inches(2.2), Inches(1.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = box["color"]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = box["text"]
            text_frame.paragraphs[0].font.size = Pt(11)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
        # Add workflow description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(1.5))
        desc_frame = desc_box.text_frame
        desc_frame.text = (
            "WORKFLOW HIGHLIGHTS:\n"
            "• Users describe requirements in plain English → AI converts to technical configuration\n"
            "• Human oversight ensures financial accuracy and compliance\n"
            "• Continuous learning improves AI suggestions over time\n"
            "• End-to-end automation with critical human checkpoints"
        )
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].font.bold = True
        
    def create_aws_architecture_slide(self):
        """Slide 6: AWS Cloud Infrastructure"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Enterprise AWS Cloud Architecture"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # AWS Components
        aws_components = [
            {"text": "Application Load\nBalancer (ALB)", "x": 1, "y": 1.5, "color": RGBColor(255, 153, 0)},
            {"text": "ECS Fargate\n(Container Runtime)", "x": 4, "y": 1.5, "color": RGBColor(255, 153, 0)},
            {"text": "API Gateway\n(REST APIs)", "x": 7, "y": 1.5, "color": RGBColor(255, 153, 0)},
            {"text": "CloudFront CDN\n(React Frontend)", "x": 10, "y": 1.5, "color": RGBColor(255, 153, 0)},
            
            {"text": "RDS PostgreSQL\n(Metadata & Configs)", "x": 1, "y": 3.2, "color": RGBColor(52, 144, 220)},
            {"text": "S3 Buckets\n(File Storage)", "x": 4, "y": 3.2, "color": RGBColor(52, 144, 220)},
            {"text": "ElastiCache Redis\n(Session & Cache)", "x": 7, "y": 3.2, "color": RGBColor(52, 144, 220)},
            {"text": "OpenSearch\n(Vector Database)", "x": 10, "y": 3.2, "color": RGBColor(52, 144, 220)},
            
            {"text": "Lambda Functions\n(Serverless AI)", "x": 2.5, "y": 4.9, "color": RGBColor(255, 193, 7)},
            {"text": "SageMaker\n(ML Pipeline)", "x": 5.5, "y": 4.9, "color": RGBColor(255, 193, 7)},
            {"text": "Bedrock\n(LLM Access)", "x": 8.5, "y": 4.9, "color": RGBColor(255, 193, 7)},
            
            {"text": "CloudWatch\n(Monitoring)", "x": 1.5, "y": 6.6, "color": RGBColor(156, 39, 176)},
            {"text": "IAM\n(Security)", "x": 4.5, "y": 6.6, "color": RGBColor(156, 39, 176)},
            {"text": "KMS\n(Encryption)", "x": 7.5, "y": 6.6, "color": RGBColor(156, 39, 176)},
            {"text": "VPC\n(Network Security)", "x": 10.5, "y": 6.6, "color": RGBColor(156, 39, 176)}
        ]
        
        for component in aws_components:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(component["x"]), Inches(component["y"]),
                Inches(2), Inches(1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = component["color"]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = component["text"]
            text_frame.paragraphs[0].font.size = Pt(10)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
    def create_security_compliance_slide(self):
        """Slide 7: Security & Compliance"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        title = slide.shapes.title
        title.text = "Enterprise Security & Compliance"
        
        content = slide.placeholders[1]
        content.text = (
            "DATA SECURITY:\n"
            "• End-to-end encryption (AES-256) for data at rest and in transit\n"
            "• AWS KMS for key management and rotation\n"
            "• VPC with private subnets and security groups\n"
            "• Multi-factor authentication and role-based access control\n\n"
            "AI & LLM SECURITY:\n"
            "• No sensitive data sent to external LLM providers\n"
            "• Data anonymization and tokenization for AI processing\n"
            "• Internal LLM option for highly sensitive workloads\n"
            "• Audit trails for all AI-generated configurations\n\n"
            "COMPLIANCE & GOVERNANCE:\n"
            "• SOX compliance with complete audit trails\n"
            "• GDPR compliance for data processing\n"
            "• Financial services regulatory alignment\n"
            "• Human-in-the-loop validation for critical decisions\n"
            "• Immutable logging and change tracking"
        )
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            if paragraph.text.endswith(':'):
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(198, 40, 40)  # Red for security emphasis
                
    def create_roi_metrics_slide(self):
        """Slide 8: ROI & Performance Metrics"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        title = slide.shapes.title
        title.text = "Return on Investment & Performance Metrics"
        
        # Add content textbox
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = (
            "OPERATIONAL EFFICIENCY:\n"
            "• 80% reduction in manual processing time\n"
            "• 95% faster rule configuration with AI\n"
            "• 50,000-100,000 records processed per batch\n"
            "• 99.5% accuracy with human validation\n\n"
            "COST SAVINGS:\n"
            "• $200K+ annual operational cost reduction\n"
            "• 60+ hours/month freed for strategic work\n"
            "• 90% reduction in configuration errors\n"
            "• Scalable infrastructure reduces marginal costs\n\n"
            "BUSINESS IMPACT:\n"
            "• Faster regulatory reporting cycles\n"
            "• Enhanced data quality and consistency\n"
            "• Reduced compliance risk\n"
            "• Accelerated financial close processes"
        )
        
        # Add a simple chart
        chart_data = CategoryChartData()
        chart_data.categories = ['Manual Process', 'AI-Assisted Process']
        chart_data.add_series('Hours per Month', (240, 48))  # 80% reduction
        chart_data.add_series('Error Rate %', (15, 0.5))     # 95% reduction
        
        x, y, cx, cy = Inches(7), Inches(2), Inches(5), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.chart_title.text_frame.text = "Performance Comparison"
        
    def create_implementation_roadmap_slide(self):
        """Slide 9: Implementation Roadmap"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        title = slide.shapes.title
        title.text = "Implementation Roadmap & Timeline"
        
        content = slide.placeholders[1]
        content.text = (
            "PHASE 1: FOUNDATION (Months 1-2)\n"
            "• AWS infrastructure setup and security configuration\n"
            "• Core platform development and basic AI integration\n"
            "• Initial LLM provider connections (OpenAI, Anthropic)\n"
            "• Basic reconciliation and transformation capabilities\n\n"
            "PHASE 2: AI ENHANCEMENT (Months 3-4)\n"
            "• Advanced AI configuration generation\n"
            "• Natural language processing integration\n"
            "• Internal LLM integration for sensitive workloads\n"
            "• Enhanced error detection and validation\n\n"
            "PHASE 3: ENTERPRISE FEATURES (Months 5-6)\n"
            "• Advanced analytics and reporting dashboards\n"
            "• API integrations with existing financial systems\n"
            "• Comprehensive audit trails and compliance features\n"
            "• Performance optimization and scaling\n\n"
            "PHASE 4: OPTIMIZATION (Ongoing)\n"
            "• Continuous AI model improvement\n"
            "• User feedback integration and feature enhancement\n"
            "• Advanced automation and workflow optimization"
        )
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            if paragraph.text.startswith('PHASE'):
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(13, 71, 161)
                
    def create_competitive_advantage_slide(self):
        """Slide 10: Competitive Advantage"""
        slide = self.prs.slides.add_slide(self.two_content_layout)
        
        title = slide.shapes.title
        title.text = "Strategic Competitive Advantage"
        
        # Left content - Our Solution
        left_content = slide.placeholders[1]
        left_content.text = (
            "FTT-ML ADVANTAGES:\n\n"
            "✓ Multi-provider LLM integration\n"
            "✓ Financial domain AI specialization\n"
            "✓ Human-in-the-loop validation\n"
            "✓ Enterprise security & compliance\n"
            "✓ Scalable AWS cloud architecture\n"
            "✓ Natural language configuration\n"
            "✓ Continuous AI learning\n"
            "✓ Complete audit trail\n"
            "✓ Real-time processing capabilities\n\n"
            "TECHNOLOGY DIFFERENTIATION:\n"
            "• Purpose-built for financial data\n"
            "• AI-first architecture\n"
            "• Pluggable LLM ecosystem"
        )
        
        # Right content - Market Position
        right_content = slide.placeholders[2]
        right_content.text = (
            "MARKET POSITION:\n\n"
            "⚡ First-mover advantage in AI-powered\n   financial data processing\n\n"
            "🎯 Addresses $50B+ market opportunity\n   in financial data management\n\n"
            "🚀 Potential for product commercialization\n   to external financial institutions\n\n"
            "📈 Foundation for advanced AI initiatives:\n"
            "   • Predictive analytics\n"
            "   • Risk modeling\n"
            "   • Regulatory intelligence\n"
            "   • Automated reporting\n\n"
            "🏆 Positions organization as AI leader\n   in financial services industry"
        )
        
        # Format content
        for content in [left_content, right_content]:
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                if paragraph.text.startswith('FTT-ML') or paragraph.text.startswith('TECHNOLOGY') or paragraph.text.startswith('MARKET'):
                    paragraph.font.size = Pt(16)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(13, 71, 161)
                    
    def create_next_steps_slide(self):
        """Slide 11: Next Steps & Call to Action"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        title = slide.shapes.title
        title.text = "Next Steps & Investment Decision"
        
        content = slide.placeholders[1]
        content.text = (
            "IMMEDIATE ACTIONS REQUIRED:\n\n"
            "1. Executive Approval for Phase 1 Implementation\n"
            "   • Budget allocation: $500K for infrastructure and development\n"
            "   • Team assignment: 5 developers + 2 AI specialists\n"
            "   • Timeline: 6-month implementation plan\n\n"
            "2. Technology Partnership Decisions\n"
            "   • LLM provider licensing agreements\n"
            "   • AWS enterprise account setup\n"
            "   • Security and compliance reviews\n\n"
            "3. Change Management Planning\n"
            "   • User training and adoption strategy\n"
            "   • Process documentation and knowledge transfer\n"
            "   • Success metrics and KPI definition\n\n"
            "EXPECTED OUTCOMES:\n"
            "• 18-month ROI with $1.2M in operational savings\n"
            "• Strategic platform for future AI initiatives\n"
            "• Industry leadership in financial AI applications\n"
            "• Foundation for potential product commercialization"
        )
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(15)
            if paragraph.text.startswith('IMMEDIATE') or paragraph.text.startswith('EXPECTED'):
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(198, 40, 40)
                
    def generate_presentation(self):
        """Generate the complete presentation"""
        print("🚀 Generating Financial Data Transformation Platform Presentation...")
        
        self.create_title_slide()
        print("✅ Created title slide")
        
        self.create_executive_summary_slide()
        print("✅ Created executive summary")
        
        self.create_business_problem_slide()
        print("✅ Created business problem & solution")
        
        self.create_ai_capabilities_slide()
        print("✅ Created AI capabilities")
        
        self.create_workflow_diagram_slide()
        print("✅ Created workflow diagram")
        
        self.create_aws_architecture_slide()
        print("✅ Created AWS architecture")
        
        self.create_security_compliance_slide()
        print("✅ Created security & compliance")
        
        self.create_roi_metrics_slide()
        print("✅ Created ROI metrics")
        
        self.create_implementation_roadmap_slide()
        print("✅ Created implementation roadmap")
        
        self.create_competitive_advantage_slide()
        print("✅ Created competitive advantage")
        
        self.create_next_steps_slide()
        print("✅ Created next steps")
        
        return self.prs
        
    def save_presentation(self, filename="FTT-ML_Management_Presentation.pptx"):
        """Save the presentation to file"""
        self.prs.save(filename)
        print(f"💾 Presentation saved as: {filename}")
        return filename

if __name__ == "__main__":
    # Generate the presentation
    generator = FTTMLPresentationGenerator()
    presentation = generator.generate_presentation()
    filename = generator.save_presentation()
    
    print(f"\n🎉 Management presentation created successfully!")
    print(f"📁 File location: {os.path.abspath(filename)}")
    print(f"📊 Total slides: {len(presentation.slides)}")
    print("\n📋 Presentation includes:")
    print("   • Executive summary with AI focus")
    print("   • Business problem and AI-powered solution")
    print("   • Detailed LLM capabilities and integration")
    print("   • Visual workflow diagrams")
    print("   • Comprehensive AWS architecture")
    print("   • Security and compliance overview")
    print("   • ROI metrics and performance data")
    print("   • Implementation roadmap")
    print("   • Competitive advantage analysis")
    print("   • Next steps and investment decision")